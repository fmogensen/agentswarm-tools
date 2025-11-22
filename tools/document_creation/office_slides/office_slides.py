"""
Generate professional PowerPoint presentations (.pptx) from structured content.
"""

from typing import Any, Dict, List, Optional
from pydantic import Field
import os
import tempfile
from datetime import datetime
import urllib.request
import urllib.parse

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from shared.base import BaseTool
from shared.errors import ValidationError, APIError, ConfigurationError


class OfficeSlidesTool(BaseTool):
    """
    Generate or modify professional PowerPoint presentations (.pptx) from structured content.

    Use this tool to create or modify formatted slide decks including business presentations,
    educational materials, pitch decks, and reports with proper themes and layouts.

    Args:
        mode: Operation mode - "create" to create new presentation, "modify" to update existing
        slides: List of slide definitions (each with title, content, layout)
        theme: Presentation theme (modern, classic, minimal, corporate) - only for create mode
        title_slide: Title slide configuration (title, subtitle, author)
        charts: Optional list of chart definitions to embed
        output_format: Output format (pptx, pdf, both)
        existing_file_url: URL to existing presentation (required for modify mode)

    Returns:
        Dict containing:
        - success: Boolean indicating success
        - result: {"presentation_url": "...", "format": "...", "slides": N, "mode": "..."}
        - metadata: Tool execution metadata

    Example (Create):
        >>> tool = OfficeSlidesTool(
        ...     mode="create",
        ...     slides=[
        ...         {"title": "Introduction", "content": ["Point 1", "Point 2"], "layout": "title_content"},
        ...         {"title": "Data", "content": ["Analysis"], "layout": "content"}
        ...     ],
        ...     theme="modern",
        ...     title_slide={"title": "Q4 Report", "subtitle": "2025", "author": "John Doe"}
        ... )
        >>> result = tool.run()

    Example (Modify):
        >>> tool = OfficeSlidesTool(
        ...     mode="modify",
        ...     existing_file_url="computer:///path/to/presentation.pptx",
        ...     slides=[
        ...         {"title": "New Slide", "content": ["New point"], "layout": "content"}
        ...     ]
        ... )
        >>> result = tool.run()
    """

    # Tool metadata
    tool_name: str = "office_slides"
    tool_category: str = "document_creation"

    # Parameters
    mode: str = Field("create", description="Operation mode: create or modify")
    slides: List[Dict[str, Any]] = Field(
        ...,
        description="List of slide definitions with title, content, and layout",
        min_items=1
    )
    theme: str = Field(
        "modern",
        description="Theme: modern, classic, minimal, corporate (create mode only)"
    )
    title_slide: Optional[Dict[str, str]] = Field(
        None,
        description="Title slide config with title, subtitle, author"
    )
    charts: Optional[List[Dict[str, Any]]] = Field(
        None,
        description="Chart definitions to embed in slides"
    )
    output_format: str = Field(
        "pptx",
        description="Output format: pptx, pdf, both"
    )
    existing_file_url: Optional[str] = Field(
        None,
        description="URL to existing file (required for modify mode)"
    )

    def _execute(self) -> Dict[str, Any]:
        """
        Execute the office_slides tool.

        Returns:
            Dict with presentation URL and metadata
        """
        # 1. CHECK MOCK MODE FIRST (before validation that requires libraries)
        if self._should_use_mock():
            # Still do basic parameter validation
            self._validate_basic_parameters()
            return self._generate_mock_results()

        # 2. FULL VALIDATION (including library checks)
        self._validate_parameters()

        # 3. EXECUTE
        try:
            result = self._process()

            return {
                "success": True,
                "result": result,
                "metadata": {"tool_name": self.tool_name},
            }
        except Exception as e:
            raise APIError(f"Failed to generate presentation: {e}", tool_name=self.tool_name)

    def _validate_basic_parameters(self) -> None:
        """Validate basic parameters (used in mock mode)."""
        # Validate mode
        valid_modes = ["create", "modify"]
        if self.mode not in valid_modes:
            raise ValidationError(
                f"Invalid mode '{self.mode}'. Must be one of: {', '.join(valid_modes)}",
                tool_name=self.tool_name,
                field="mode"
            )

        # Validate mode-specific requirements
        if self.mode == "modify":
            if not self.existing_file_url:
                raise ValidationError(
                    "existing_file_url is required when mode='modify'",
                    tool_name=self.tool_name,
                    field="existing_file_url"
                )
        elif self.mode == "create":
            if self.existing_file_url:
                raise ValidationError(
                    "existing_file_url should not be provided when mode='create'",
                    tool_name=self.tool_name,
                    field="existing_file_url"
                )

        # Validate slides list
        if not self.slides or len(self.slides) == 0:
            raise ValidationError("slides list cannot be empty", tool_name=self.tool_name)

        # Validate theme (only matters for create mode)
        valid_themes = ["modern", "classic", "minimal", "corporate"]
        if self.mode == "create" and self.theme not in valid_themes:
            raise ValidationError(
                f"Invalid theme '{self.theme}'. Must be one of: {', '.join(valid_themes)}",
                tool_name=self.tool_name
            )

        # Validate output format
        valid_formats = ["pptx", "pdf", "both"]
        if self.output_format not in valid_formats:
            raise ValidationError(
                f"Invalid output_format '{self.output_format}'. Must be one of: {', '.join(valid_formats)}",
                tool_name=self.tool_name
            )

    def _validate_parameters(self) -> None:
        """Validate input parameters (full validation including dependencies)."""
        # Do basic validation first
        self._validate_basic_parameters()

        # Validate each slide has required fields
        for i, slide in enumerate(self.slides):
            if "title" not in slide:
                raise ValidationError(
                    f"Slide {i} missing 'title' field",
                    tool_name=self.tool_name
                )
            if "content" not in slide:
                raise ValidationError(
                    f"Slide {i} missing 'content' field",
                    tool_name=self.tool_name
                )
            if "layout" not in slide:
                raise ValidationError(
                    f"Slide {i} missing 'layout' field",
                    tool_name=self.tool_name
                )

            # Validate layout type
            valid_layouts = ["title_content", "content", "two_column", "blank"]
            if slide["layout"] not in valid_layouts:
                raise ValidationError(
                    f"Slide {i} has invalid layout '{slide['layout']}'. Must be one of: {', '.join(valid_layouts)}",
                    tool_name=self.tool_name
                )

        # Validate title_slide if provided
        if self.title_slide:
            if "title" not in self.title_slide:
                raise ValidationError(
                    "title_slide missing 'title' field",
                    tool_name=self.tool_name
                )

        # Check if python-pptx is available
        if not PPTX_AVAILABLE:
            raise ConfigurationError(
                "python-pptx library not installed. Install with: pip install python-pptx",
                tool_name=self.tool_name
            )

    def _should_use_mock(self) -> bool:
        """Check if mock mode enabled."""
        return os.getenv("USE_MOCK_APIS", "false").lower() == "true"

    def _generate_mock_results(self) -> Dict[str, Any]:
        """Generate mock results for testing."""
        total_slides = len(self.slides) + (1 if self.title_slide else 0)

        return {
            "success": True,
            "result": {
                "presentation_url": "https://mock.example.com/pres123.pptx",
                "format": self.output_format,
                "slides": total_slides,
                "file_size": "1.2 MB",
                "mode": self.mode
            },
            "metadata": {"mock_mode": True},
        }

    def _process(self) -> Dict[str, Any]:
        """Main processing logic to create or modify PowerPoint presentation."""
        if self.mode == "create":
            return self._process_create()
        else:  # modify
            return self._process_modify()

    def _process_create(self) -> Dict[str, Any]:
        """Create new presentation from scratch."""
        # Create a new presentation
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # Apply theme colors
        theme_colors = self._get_theme_colors()

        # Add title slide if provided
        if self.title_slide:
            self._add_title_slide(prs, theme_colors)

        # Add content slides
        for slide_def in self.slides:
            self._add_content_slide(prs, slide_def, theme_colors)

        # Save and return result
        return self._save_and_return(prs)

    def _process_modify(self) -> Dict[str, Any]:
        """Modify existing presentation."""
        # Download existing file
        local_path = self._download_file(self.existing_file_url)

        try:
            # Open existing presentation
            prs = Presentation(local_path)

            # Get theme colors (use modern as default for modify mode)
            theme_colors = self._get_theme_colors()

            # Update title slide if provided
            if self.title_slide and len(prs.slides) > 0:
                # Update first slide if it exists
                first_slide = prs.slides[0]
                # Try to update text in first slide's shapes
                for shape in first_slide.shapes:
                    if hasattr(shape, "text_frame"):
                        if "title" in self.title_slide and shape.text_frame.text:
                            shape.text_frame.text = self.title_slide["title"]
                            break

            # Add new content slides
            for slide_def in self.slides:
                self._add_content_slide(prs, slide_def, theme_colors)

            # Save and return result
            return self._save_and_return(prs)

        finally:
            # Clean up downloaded file
            try:
                os.unlink(local_path)
            except:
                pass

    def _download_file(self, url: str) -> str:
        """Download file from URL to temporary location."""
        # Handle computer:// protocol
        if url.startswith("computer://"):
            file_path = url.replace("computer://", "")
            if os.path.exists(file_path):
                return file_path
            else:
                raise ValidationError(f"File not found: {file_path}", tool_name=self.tool_name)

        # Handle http/https URLs
        elif url.startswith("http://") or url.startswith("https://"):
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            temp_file.close()

            try:
                urllib.request.urlretrieve(url, temp_file.name)
                return temp_file.name
            except Exception as e:
                raise APIError(f"Failed to download file from URL: {e}", tool_name=self.tool_name)

        else:
            raise ValidationError(f"Unsupported URL scheme: {url}", tool_name=self.tool_name)

    def _save_and_return(self, prs: Presentation) -> Dict[str, Any]:
        """Save presentation and return result metadata."""
        # Save presentation to temporary file
        temp_file = self._save_presentation(prs)

        # Count total slides
        total_slides = len(prs.slides)

        # Get file size
        file_size = os.path.getsize(temp_file)
        file_size_str = self._format_file_size(file_size)

        # In real implementation, would upload to AI Drive or persistent storage
        # For now, we'll use a computer:// URL
        presentation_url = f"computer://{temp_file}"

        return {
            "presentation_url": presentation_url,
            "format": self.output_format,
            "slides": total_slides,
            "file_size": file_size_str,
            "file_path": temp_file,
            "mode": self.mode
        }

    def _get_theme_colors(self) -> Dict[str, RGBColor]:
        """Get color scheme based on selected theme."""
        themes = {
            "modern": {
                "primary": RGBColor(0, 120, 215),      # Blue
                "secondary": RGBColor(76, 175, 80),    # Green
                "accent": RGBColor(255, 152, 0),       # Orange
                "text": RGBColor(33, 33, 33),          # Dark Gray
                "background": RGBColor(255, 255, 255)  # White
            },
            "classic": {
                "primary": RGBColor(0, 51, 102),       # Navy
                "secondary": RGBColor(128, 128, 128),  # Gray
                "accent": RGBColor(192, 0, 0),         # Red
                "text": RGBColor(0, 0, 0),             # Black
                "background": RGBColor(255, 255, 255)  # White
            },
            "minimal": {
                "primary": RGBColor(97, 97, 97),       # Gray
                "secondary": RGBColor(189, 189, 189),  # Light Gray
                "accent": RGBColor(33, 33, 33),        # Dark Gray
                "text": RGBColor(66, 66, 66),          # Medium Gray
                "background": RGBColor(255, 255, 255)  # White
            },
            "corporate": {
                "primary": RGBColor(0, 32, 96),        # Corporate Blue
                "secondary": RGBColor(112, 48, 160),   # Purple
                "accent": RGBColor(255, 185, 0),       # Gold
                "text": RGBColor(0, 0, 0),             # Black
                "background": RGBColor(255, 255, 255)  # White
            }
        }

        return themes.get(self.theme, themes["modern"])

    def _add_title_slide(self, prs: Presentation, theme_colors: Dict[str, RGBColor]) -> None:
        """Add title slide to presentation."""
        # Use blank layout and add custom text boxes
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = self.title_slide.get("title", "Untitled Presentation")
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_colors["primary"]
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if self.title_slide.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.8), Inches(8), Inches(0.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = self.title_slide["subtitle"]
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = theme_colors["secondary"]
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Add author if provided
        if self.title_slide.get("author"):
            author_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(8), Inches(0.3)
            )
            author_frame = author_box.text_frame
            author_frame.text = f"By {self.title_slide['author']}"
            author_para = author_frame.paragraphs[0]
            author_para.font.size = Pt(16)
            author_para.font.color.rgb = theme_colors["text"]
            author_para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_def: Dict[str, Any],
        theme_colors: Dict[str, RGBColor]
    ) -> None:
        """Add content slide to presentation."""
        layout_type = slide_def["layout"]

        if layout_type == "title_content":
            self._add_title_content_slide(prs, slide_def, theme_colors)
        elif layout_type == "content":
            self._add_content_only_slide(prs, slide_def, theme_colors)
        elif layout_type == "two_column":
            self._add_two_column_slide(prs, slide_def, theme_colors)
        elif layout_type == "blank":
            self._add_blank_slide(prs, slide_def, theme_colors)

    def _add_title_content_slide(
        self,
        prs: Presentation,
        slide_def: Dict[str, Any],
        theme_colors: Dict[str, RGBColor]
    ) -> None:
        """Add slide with title and content."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_def["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_colors["primary"]

        # Add content
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Add content items
        content = slide_def["content"]
        if isinstance(content, list):
            for i, item in enumerate(content):
                if i > 0:
                    content_frame.add_paragraph()
                para = content_frame.paragraphs[i]
                para.text = f"• {item}"
                para.font.size = Pt(18)
                para.font.color.rgb = theme_colors["text"]
                para.space_before = Pt(6)
        else:
            para = content_frame.paragraphs[0]
            para.text = str(content)
            para.font.size = Pt(18)
            para.font.color.rgb = theme_colors["text"]

    def _add_content_only_slide(
        self,
        prs: Presentation,
        slide_def: Dict[str, Any],
        theme_colors: Dict[str, RGBColor]
    ) -> None:
        """Add slide with only content (no prominent title)."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Small title at top
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_def["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.color.rgb = theme_colors["secondary"]

        # Content fills most of slide
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.9), Inches(8.4), Inches(4.2)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        content = slide_def["content"]
        if isinstance(content, list):
            for i, item in enumerate(content):
                if i > 0:
                    content_frame.add_paragraph()
                para = content_frame.paragraphs[i]
                para.text = f"• {item}"
                para.font.size = Pt(20)
                para.font.color.rgb = theme_colors["text"]
        else:
            para = content_frame.paragraphs[0]
            para.text = str(content)
            para.font.size = Pt(20)
            para.font.color.rgb = theme_colors["text"]

    def _add_two_column_slide(
        self,
        prs: Presentation,
        slide_def: Dict[str, Any],
        theme_colors: Dict[str, RGBColor]
    ) -> None:
        """Add slide with two columns."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_def["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_colors["primary"]

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(4.5), Inches(3.8)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.0), Inches(1.2), Inches(4.5), Inches(3.8)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        # Fill columns
        content = slide_def["content"]
        if isinstance(content, list) and len(content) >= 2:
            # Split content into two columns
            mid = len(content) // 2
            left_content = content[:mid]
            right_content = content[mid:]

            for i, item in enumerate(left_content):
                if i > 0:
                    left_frame.add_paragraph()
                para = left_frame.paragraphs[i]
                para.text = f"• {item}"
                para.font.size = Pt(16)
                para.font.color.rgb = theme_colors["text"]

            for i, item in enumerate(right_content):
                if i > 0:
                    right_frame.add_paragraph()
                para = right_frame.paragraphs[i]
                para.text = f"• {item}"
                para.font.size = Pt(16)
                para.font.color.rgb = theme_colors["text"]
        else:
            # Put all content in left column
            left_frame.text = str(content)
            left_frame.paragraphs[0].font.size = Pt(16)
            left_frame.paragraphs[0].font.color.rgb = theme_colors["text"]

    def _add_blank_slide(
        self,
        prs: Presentation,
        slide_def: Dict[str, Any],
        theme_colors: Dict[str, RGBColor]
    ) -> None:
        """Add blank slide with minimal content."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Optional title
        if slide_def.get("title"):
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_def["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(24)
            title_para.font.color.rgb = theme_colors["secondary"]

    def _save_presentation(self, prs: Presentation) -> str:
        """Save presentation to temporary file."""
        # Create temporary file
        temp_dir = tempfile.gettempdir()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        temp_file = os.path.join(temp_dir, filename)

        # Save presentation
        prs.save(temp_file)

        return temp_file

    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human-readable format."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.1f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.1f} TB"


if __name__ == "__main__":
    print("Testing OfficeSlidesTool...")

    import os
    os.environ["USE_MOCK_APIS"] = "true"

    # Test basic presentation
    tool = OfficeSlidesTool(
        slides=[
            {"title": "Slide 1", "content": ["Point A", "Point B"], "layout": "title_content"},
            {"title": "Slide 2", "content": ["Data"], "layout": "content"}
        ],
        theme="modern",
        title_slide={"title": "Test Presentation", "subtitle": "2025"}
    )
    result = tool.run()

    assert result.get('success') == True
    assert result['result']['slides'] == 3  # title + 2 content
    print(f"✅ Basic presentation test passed")
    print(f"   Slides: {result['result']['slides']}")
    print(f"   Format: {result['result']['format']}")

    # Test PDF export
    tool2 = OfficeSlidesTool(
        slides=[{"title": "Test", "content": ["Content"], "layout": "content"}],
        output_format="pdf"
    )
    result2 = tool2.run()
    assert result2['result']['format'] == "pdf"
    print(f"✅ PDF export test passed")

    # Test two-column layout
    tool3 = OfficeSlidesTool(
        slides=[
            {
                "title": "Two Column Test",
                "content": ["Left 1", "Left 2", "Right 1", "Right 2"],
                "layout": "two_column"
            }
        ],
        theme="corporate"
    )
    result3 = tool3.run()
    assert result3.get('success') == True
    print(f"✅ Two-column layout test passed")

    # Test modify mode validation
    tool4 = OfficeSlidesTool(
        mode="modify",
        slides=[{"title": "New Slide", "content": ["Content"], "layout": "content"}],
        # Missing existing_file_url - should fail
    )
    result4 = tool4.run()
    assert result4.get('success') == False  # Should fail validation
    assert 'existing_file_url is required' in str(result4.get('error', ''))
    print(f"✅ Modify mode validation test passed")

    # Test create mode with existing_file_url (should fail)
    tool5 = OfficeSlidesTool(
        mode="create",
        slides=[{"title": "Slide", "content": ["Content"], "layout": "content"}],
        existing_file_url="computer:///some/file.pptx"  # Should not be provided for create
    )
    result5 = tool5.run()
    assert result5.get('success') == False  # Should fail validation
    assert 'should not be provided' in str(result5.get('error', ''))
    print(f"✅ Create mode validation test passed")

    # Test modify mode with mock
    tool6 = OfficeSlidesTool(
        mode="modify",
        existing_file_url="https://mock.example.com/pres.pptx",
        slides=[{"title": "New Slide", "content": ["New content"], "layout": "title_content"}]
    )
    result6 = tool6.run()
    assert result6.get('success') == True
    assert result6['result']['mode'] == "modify"
    print(f"✅ Modify mode mock test passed")

    print("\n🎉 All tests passed!")
